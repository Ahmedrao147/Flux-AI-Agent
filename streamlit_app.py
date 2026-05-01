"""
Bob AI Agent with Groq API Integration
Enhanced version with AI-powered note generation using Groq's LLM
"""

import os
import re
from typing import List, Dict, Optional
from pathlib import Path
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# File processing imports
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract

# Text processing
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

# PDF generation
from fpdf import FPDF

# Groq API
from groq import Groq
import httpx


# If you actually need a proxy, uncomment and configure this:
# client = Groq(
#     api_key="your_api_key",
#     http_client=httpx.Client(proxy="http://your-actual-proxy-url:8080")
# )

class BobAIAgentGroq:
    """Enhanced Bob AI Agent with Groq API integration."""
    
    def __init__(self, groq_api_key: str = None):
        """
        Initialize the Bob AI Agent with Groq.
        
        Args:
            groq_api_key: Groq API key (if not provided, uses environment variable from .env)
        """
        self.supported_formats = ['.pdf', '.ppt', '.pptx', '.txt', '.png', '.jpg', '.jpeg']
        self._setup_nltk()
        
        # Initialize Groq client - loads from .env file
        api_key = groq_api_key or os.getenv('GROQ_API_KEY')
        if not api_key:
            raise ValueError("GROQ_API_KEY not found. Please set it in .env file or pass it as parameter.")
        
        self.groq_client = Groq(api_key=api_key)
        self.model = "mixtral-8x7b-32768"  # Fast and capable model
    
    def _setup_nltk(self):
        """Download required NLTK data."""
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
        
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords', quiet=True)
    
    def process_file(self, file_path: str) -> Dict[str, any]:
        """
        Process input file and extract text content.
        
        Args:
            file_path: Path to the input file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {extension}")
        
        # Extract text based on file type
        if extension == '.pdf':
            text = self._extract_from_pdf(file_path)
        elif extension in ['.ppt', '.pptx']:
            text = self._extract_from_ppt(file_path)
        elif extension == '.txt':
            text = self._extract_from_txt(file_path)
        elif extension in ['.png', '.jpg', '.jpeg']:
            text = self._extract_from_image(file_path)
        else:
            text = ""
        
        return {
            'text': text,
            'filename': file_path.name,
            'format': extension
        }
    
    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error extracting PDF: {e}")
        return text
    
    def _extract_from_ppt(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting PPT: {e}")
        return text
    
    def _extract_from_txt(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading TXT: {e}")
            return ""
    
    def _extract_from_image(self, file_path: Path) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            print(f"Error performing OCR: {e}")
            return ""
    
    def create_notes_with_ai(self, text: str, title: Optional[str] = None) -> Dict[str, any]:
        """
        Convert extracted text into structured notes using Groq AI.
        
        Args:
            text: Raw text content
            title: Optional title for the notes
            
        Returns:
            Dictionary containing structured notes
        """
        if not text or not text.strip():
            return {
                'title': title or "Empty Document",
                'sections': [],
                'key_concepts': [],
                'summary': "No content to summarize."
            }
        
        # Truncate text if too long (to fit in context window)
        max_chars = 25000
        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n[Content truncated due to length...]"
        
        # Create prompt for Groq
        prompt = self._create_note_generation_prompt(text, title)
        
        try:
            # Call Groq API
            response = self.groq_client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are Bob AI, an expert note-taking assistant for students. You create clear, concise, and well-structured notes from lecture materials. Always follow the exact format specified in the user's instructions."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                temperature=0.3,
                max_tokens=4000,
            )
            
            # Parse AI response
            ai_notes = response.choices[0].message.content
            
            # Validate response before parsing
            if not ai_notes or len(ai_notes.strip()) < 50:
                print("Warning: AI response too short, using fallback")
                return self._create_notes_fallback(text, title)
            
            structured_notes = self._parse_ai_response(ai_notes, title)
            
            # Validate parsed notes
            if not self._validate_notes(structured_notes):
                print("Warning: AI response validation failed, using fallback")
                return self._create_notes_fallback(text, title)
            
            return structured_notes
            
        except Exception as e:
            print(f"Error calling Groq API: {e}")
            # Fallback to basic processing
            return self._create_notes_fallback(text, title)
    
    def _validate_notes(self, notes: Dict[str, any]) -> bool:
        """Validate that notes have minimum required content."""
        try:
            # Check title exists and is not empty
            if not notes.get('title') or len(notes['title'].strip()) < 3:
                return False
            
            # Check we have at least some sections or key concepts
            has_sections = notes.get('sections') and len(notes['sections']) > 0
            has_concepts = notes.get('key_concepts') and len(notes['key_concepts']) > 0
            
            if not (has_sections or has_concepts):
                return False
            
            # If we have sections, check they have content
            if has_sections:
                for section in notes['sections']:
                    if not section.get('heading') or not section.get('content'):
                        return False
                    if len(section['content']) == 0:
                        return False
            
            return True
            
        except Exception as e:
            print(f"Validation error: {e}")
            return False
    
    def _create_note_generation_prompt(self, text: str, title: Optional[str]) -> str:
        """Create a prompt for the AI to generate notes."""
        title_instruction = f'use the title "{title}"' if title else "generate an appropriate title"
        
        prompt = f"""You are Bob AI, an expert note-taking assistant. Convert the following lecture content into structured, student-friendly notes.

LECTURE CONTENT:
{text}

CRITICAL FORMATTING INSTRUCTIONS - YOU MUST FOLLOW THIS EXACT FORMAT:

1. Start with "TITLE:" followed by the title ({title_instruction})
2. Add "KEY CONCEPTS:" section with 8-10 bullet points (use • or -)
3. Create 3-5 "SECTION:" blocks, each with a descriptive heading
4. Each section should have 3-8 bullet points (use • or -)
5. End with "SUMMARY:" followed by 2-3 sentences
6. Use clear, simple language suitable for students
7. Focus on the most important information

REQUIRED FORMAT (DO NOT DEVIATE):

TITLE: [Your title here]

KEY CONCEPTS:
• First key concept or term
• Second key concept or term
• Third key concept or term
• [Continue with 5-7 more concepts]

SECTION: [First Section Heading]
• First important point about this topic
• Second important point with details
• Third point explaining the concept
• [Add 2-5 more relevant points]

SECTION: [Second Section Heading]
• First point for this section
• Second point with explanation
• Third point with examples
• [Add more points as needed]

SECTION: [Third Section Heading]
• Key point one
• Key point two
• Key point three
• [Continue as needed]

SUMMARY:
[Write 2-3 clear sentences summarizing the entire content. Make it concise and informative.]

IMPORTANT REMINDERS:
- Use EXACTLY the labels: TITLE:, KEY CONCEPTS:, SECTION:, SUMMARY:
- Start each bullet point with • or -
- Keep each bullet point to 1-2 sentences maximum
- Make content clear and easy to understand
- Ensure all sections have substantive content

Now generate the structured notes:"""
        
        return prompt
    
    def _parse_ai_response(self, ai_response: str, default_title: Optional[str]) -> Dict[str, any]:
        """Parse the AI-generated notes into structured format with robust error handling."""
        lines = ai_response.strip().split('\n')
        
        title = default_title or "Lecture Notes"
        key_concepts = []
        sections = []
        summary = ""
        
        current_section = None
        current_content = []
        in_summary = False
        in_concepts = False
        
        # Track if we found any structured content
        found_structure = False
        
        for line in lines:
            line = line.strip()
            
            if not line:
                continue
            
            # Extract title (case-insensitive, flexible matching)
            if line.upper().startswith('TITLE:') or line.upper().startswith('# '):
                title = line.split(':', 1)[-1].strip().lstrip('#').strip()
                if title:
                    found_structure = True
            
            # Key concepts section (flexible matching)
            elif line.upper().startswith('KEY CONCEPTS:') or line.upper().startswith('KEY TERMS:'):
                in_concepts = True
                in_summary = False
                found_structure = True
            
            # Section headers (flexible matching)
            elif line.upper().startswith('SECTION:') or (line.upper().startswith('##') and not in_summary):
                # Save previous section
                if current_section and current_content:
                    sections.append({
                        'heading': current_section,
                        'content': current_content
                    })
                # Extract section name
                if ':' in line:
                    current_section = line.split(':', 1)[1].strip()
                else:
                    current_section = line.lstrip('#').strip()
                current_content = []
                in_concepts = False
                in_summary = False
                found_structure = True
            
            # Summary section (flexible matching)
            elif line.upper().startswith('SUMMARY:') or line.upper().startswith('CONCLUSION:'):
                in_summary = True
                in_concepts = False
                # Save last section
                if current_section and current_content:
                    sections.append({
                        'heading': current_section,
                        'content': current_content
                    })
                    current_section = None
                # Extract summary text if on same line
                if ':' in line:
                    summary_start = line.split(':', 1)[1].strip()
                    if summary_start:
                        summary = summary_start + " "
                found_structure = True
            
            # Content lines (bullet points with various markers)
            elif line.startswith('•') or line.startswith('-') or line.startswith('*') or line.startswith('→'):
                content = line.lstrip('•-*→').strip()
                if content:
                    if in_concepts:
                        key_concepts.append(content)
                    elif current_section:
                        current_content.append(content)
                    found_structure = True
            
            # Numbered lists
            elif re.match(r'^\d+[\.)]\s+', line):
                content = re.sub(r'^\d+[\.)]\s+', '', line).strip()
                if content:
                    if in_concepts:
                        key_concepts.append(content)
                    elif current_section:
                        current_content.append(content)
                    found_structure = True
            
            # Summary content (multi-line)
            elif in_summary and line:
                summary += line + " "
            
            # If no structure found yet, try to extract content anyway
            elif not found_structure and len(line) > 20:
                # Treat as potential section content
                if not current_section:
                    current_section = "Main Content"
                current_content.append(line)
        
        # Save last section if exists
        if current_section and current_content:
            sections.append({
                'heading': current_section,
                'content': current_content
            })
        
        # Validation: Ensure we have minimum content
        if not found_structure or (not sections and not key_concepts):
            # Fallback: treat entire response as content
            return self._parse_unstructured_response(ai_response, default_title)
        
        # Ensure we have at least some key concepts
        if not key_concepts and sections:
            # Extract first few points from first section as concepts
            if sections[0]['content']:
                key_concepts = sections[0]['content'][:min(5, len(sections[0]['content']))]
        
        # Ensure we have a summary
        if not summary.strip() and sections:
            # Create a basic summary from section headings
            summary = f"This document covers {len(sections)} main topics: " + ", ".join([s['heading'] for s in sections[:3]])
            if len(sections) > 3:
                summary += ", and more."
        
        return {
            'title': title,
            'key_concepts': key_concepts[:15],  # Limit to 15 concepts
            'sections': sections,
            'summary': summary.strip()
        }
    
    def _parse_unstructured_response(self, text: str, default_title: Optional[str]) -> Dict[str, any]:
        """Fallback parser for unstructured AI responses."""
        # Split into paragraphs
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        # Use first paragraph as title if no default
        title = default_title or (paragraphs[0][:50] + "..." if paragraphs else "Lecture Notes")
        
        # Extract sentences for content
        sentences = []
        for para in paragraphs:
            para_sentences = sent_tokenize(para)
            sentences.extend([s.strip() for s in para_sentences if len(s.strip()) > 20])
        
        # Create sections from paragraphs
        sections = []
        for i, para in enumerate(paragraphs[:5], 1):  # Max 5 sections
            para_sentences = sent_tokenize(para)
            if para_sentences:
                sections.append({
                    'heading': f"Section {i}",
                    'content': [s.strip() for s in para_sentences if len(s.strip()) > 20]
                })
        
        # Extract key concepts using word frequency
        key_concepts = self._extract_key_concepts(text)[:10]
        
        # Create summary from first few sentences
        summary = " ".join(sentences[:3]) if len(sentences) >= 3 else " ".join(sentences)
        
        return {
            'title': title,
            'key_concepts': key_concepts,
            'sections': sections,
            'summary': summary
        }
    
    def _create_notes_fallback(self, text: str, title: Optional[str]) -> Dict[str, any]:
        """Fallback method using basic NLP if AI fails."""
        cleaned_text = self._clean_text(text)
        sentences = sent_tokenize(cleaned_text)
        sections = self._identify_sections(cleaned_text)
        summary = self._generate_summary(sentences)
        key_concepts = self._extract_key_concepts(cleaned_text)
        
        return {
            'title': title or self._generate_title(cleaned_text),
            'sections': sections,
            'key_concepts': key_concepts,
            'summary': summary
        }
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s.,!?;:()\-]', '', text)
        return text.strip()
    
    def _identify_sections(self, text: str) -> List[Dict[str, any]]:
        """Identify and structure sections in the text."""
        sections = []
        lines = text.split('\n')
        current_section = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if self._is_heading(line):
                if current_section:
                    sections.append({
                        'heading': current_section,
                        'content': self._format_content(current_content)
                    })
                current_section = line
                current_content = []
            else:
                current_content.append(line)
        
        if current_section:
            sections.append({
                'heading': current_section,
                'content': self._format_content(current_content)
            })
        
        if not sections:
            sections.append({
                'heading': 'Main Content',
                'content': self._format_content([text])
            })
        
        return sections
    
    def _is_heading(self, line: str) -> bool:
        """Determine if a line is likely a heading."""
        if len(line) < 5 or len(line) > 100:
            return False
        if re.match(r'^\d+\.(\d+\.)*\s+', line):
            return True
        if line.isupper() or line.istitle():
            return True
        if line.endswith(':'):
            return True
        return False
    
    def _format_content(self, content_lines: List[str]) -> List[str]:
        """Format content into bullet points."""
        formatted = []
        for line in content_lines:
            if line.strip():
                sentences = sent_tokenize(line)
                for sentence in sentences:
                    if len(sentence.split()) > 5:
                        formatted.append(sentence.strip())
        return formatted
    
    def _generate_summary(self, sentences: List[str], max_sentences: int = 3) -> str:
        """Generate a brief summary of the content."""
        if not sentences:
            return "No content to summarize."
        
        word_freq = {}
        stop_words = set(stopwords.words('english'))
        
        for sentence in sentences:
            words = word_tokenize(sentence.lower())
            for word in words:
                if word.isalnum() and word not in stop_words:
                    word_freq[word] = word_freq.get(word, 0) + 1
        
        sentence_scores = {}
        for sentence in sentences:
            words = word_tokenize(sentence.lower())
            score = sum(word_freq.get(word, 0) for word in words if word.isalnum())
            sentence_scores[sentence] = score
        
        top_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)
        summary_sentences = [s[0] for s in top_sentences[:max_sentences]]
        
        return ' '.join(summary_sentences)
    
    def _extract_key_concepts(self, text: str) -> List[str]:
        """Extract key concepts and terms from the text."""
        words = word_tokenize(text.lower())
        stop_words = set(stopwords.words('english'))
        
        word_freq = {}
        for word in words:
            if word.isalnum() and len(word) > 3 and word not in stop_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        top_concepts = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
        return [concept[0].capitalize() for concept in top_concepts]
    
    def _generate_title(self, text: str) -> str:
        """Generate a title from the content."""
        sentences = sent_tokenize(text)
        if sentences:
            first_sentence = sentences[0]
            words = first_sentence.split()[:6]
            return ' '.join(words) + ('...' if len(first_sentence.split()) > 6 else '')
        return "Lecture Notes"
    
    def export_to_txt(self, notes: Dict[str, any], output_path: str):
        """Export notes to TXT file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"{'='*60}\n")
            f.write(f"{notes['title'].upper()}\n")
            f.write(f"{'='*60}\n\n")
            
            if notes.get('key_concepts'):
                f.write("KEY CONCEPTS:\n")
                f.write(f"{'-'*60}\n")
                for concept in notes['key_concepts']:
                    f.write(f"• {concept}\n")
                f.write("\n")
            
            for section in notes['sections']:
                f.write(f"\n{section['heading'].upper()}\n")
                f.write(f"{'-'*60}\n")
                for point in section['content']:
                    f.write(f"• {point}\n")
                f.write("\n")
            
            if notes.get('summary'):
                f.write(f"\nSUMMARY\n")
                f.write(f"{'-'*60}\n")
                f.write(f"{notes['summary']}\n")
    
    def export_to_pdf(self, notes: Dict[str, any], output_path: str):
        """Export notes to PDF file with enhanced formatting."""
        pdf = FPDF()
        # Set margins BEFORE adding page
        pdf.set_margins(left=20, top=20, right=20)
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()
        
        # Helper function to clean text for PDF
        def clean_text(text):
            """Remove problematic characters for PDF encoding."""
            # Replace common problematic characters
            replacements = {
                '\u2018': "'", '\u2019': "'",  # Smart quotes
                '\u201c': '"', '\u201d': '"',  # Smart double quotes
                '\u2013': '-', '\u2014': '-',  # En/em dashes
                '\u2022': '-',  # Bullet point
                '\u2026': '...',  # Ellipsis
            }
            for old, new in replacements.items():
                text = text.replace(old, new)
            # Encode to latin-1 and decode, replacing errors
            return text.encode('latin-1', 'replace').decode('latin-1')
        
        # Title with decorative line
        pdf.set_font('Arial', 'B', 18)
        title_text = clean_text(notes['title'])
        pdf.cell(0, 12, title_text, ln=True, align='C')
        pdf.set_draw_color(30, 136, 229)  # Blue color
        pdf.set_line_width(0.5)
        pdf.line(20, pdf.get_y(), 190, pdf.get_y())
        pdf.ln(8)
        
        # AI Badge
        pdf.set_font('Arial', 'I', 9)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 5, 'Generated by Bob AI Agent', ln=True, align='C')
        pdf.set_text_color(0, 0, 0)
        pdf.ln(5)
        
        # Key Concepts Section
        if notes.get('key_concepts') and len(notes['key_concepts']) > 0:
            pdf.set_font('Arial', 'B', 14)
            pdf.set_fill_color(227, 242, 253)  # Light blue background
            pdf.cell(0, 10, 'KEY CONCEPTS', ln=True, fill=True)
            pdf.ln(2)
            
            pdf.set_font('Arial', '', 10)
            for i, concept in enumerate(notes['key_concepts'], 1):
                concept_text = clean_text(concept)
                # Add bullet point with indentation
                pdf.set_x(25)
                pdf.multi_cell(0, 6, f'{i}. {concept_text}', align='L')
            pdf.ln(5)
        
        # Content Sections
        for idx, section in enumerate(notes['sections']):
            # Section heading with background
            pdf.set_font('Arial', 'B', 13)
            pdf.set_fill_color(240, 240, 240)  # Light gray background
            section_heading = clean_text(section['heading'])
            pdf.cell(0, 9, section_heading, ln=True, fill=True)
            pdf.ln(2)
            
            # Section content
            pdf.set_font('Arial', '', 10)
            for point in section['content']:
                point_text = clean_text(point)
                # Add bullet point with indentation
                pdf.set_x(25)
                pdf.multi_cell(0, 6, f'- {point_text}', align='L')
            pdf.ln(4)
        
        # Summary Section
        if notes.get('summary') and notes['summary'].strip():
            pdf.set_font('Arial', 'B', 14)
            pdf.set_fill_color(255, 248, 225)  # Light yellow background
            pdf.cell(0, 10, 'SUMMARY', ln=True, fill=True)
            pdf.ln(2)
            
            pdf.set_font('Arial', '', 10)
            summary_text = clean_text(notes['summary'])
            pdf.multi_cell(0, 6, summary_text, align='L')
            pdf.ln(5)
        
        # Footer
        pdf.set_y(-30)
        pdf.set_font('Arial', 'I', 8)
        pdf.set_text_color(150, 150, 150)
        pdf.cell(0, 10, 'Created with Bob AI Agent - Your Intelligent Note-Taking Assistant', ln=True, align='C')
        
        pdf.output(output_path)


def main():
    """Main function for command-line usage."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python bob_agent_groq.py <input_file> [output_format]")
        print("Supported formats: pdf, ppt, pptx, txt, png, jpg, jpeg")
        print("Output formats: txt, pdf (default: both)")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_format = sys.argv[2] if len(sys.argv) > 2 else 'both'
    
    bob = BobAIAgentGroq()
    
    print(f"Processing file: {input_file}")
    extracted = bob.process_file(input_file)
    print(f"Extracted {len(extracted['text'])} characters")
    
    print("Generating AI-powered notes...")
    notes = bob.create_notes_with_ai(extracted['text'])
    print(f"Generated notes with {len(notes['sections'])} sections")
    
    base_name = Path(input_file).stem
    
    if output_format in ['txt', 'both']:
        txt_output = f"{base_name}_notes.txt"
        bob.export_to_txt(notes, txt_output)
        print(f"Exported to: {txt_output}")
    
    if output_format in ['pdf', 'both']:
        pdf_output = f"{base_name}_notes.pdf"
        bob.export_to_pdf(notes, pdf_output)
        print(f"Exported to: {pdf_output}")
    
    print("\n✅ AI-powered notes created successfully!")


if __name__ == "__main__":
    main()

# Made with Bob
