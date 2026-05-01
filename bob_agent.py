"""
Bob AI Agent - Intelligent Note-Taking Assistant for Students
Accepts lecture inputs in multiple formats and converts them into structured notes.
"""

import os
import re
from typing import List, Dict, Optional
from pathlib import Path

# File processing imports
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract

# Text processing
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

# PDF generation
from fpdf import FPDF


class BobAIAgent:
    """Main class for the Bob AI note-taking assistant."""
    
    def __init__(self):
        """Initialize the Bob AI Agent."""
        self.supported_formats = ['.pdf', '.ppt', '.pptx', '.txt', '.png', '.jpg', '.jpeg']
        self._setup_nltk()
    
    def _setup_nltk(self):
        """Download required NLTK data."""
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
        
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords', quiet=True)
    
    def process_file(self, file_path: str) -> Dict[str, any]:
        """
        Process input file and extract text content.
        
        Args:
            file_path: Path to the input file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {extension}")
        
        # Extract text based on file type
        if extension == '.pdf':
            text = self._extract_from_pdf(file_path)
        elif extension in ['.ppt', '.pptx']:
            text = self._extract_from_ppt(file_path)
        elif extension == '.txt':
            text = self._extract_from_txt(file_path)
        elif extension in ['.png', '.jpg', '.jpeg']:
            text = self._extract_from_image(file_path)
        else:
            text = ""
        
        return {
            'text': text,
            'filename': file_path.name,
            'format': extension
        }
    
    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error extracting PDF: {e}")
        return text
    
    def _extract_from_ppt(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting PPT: {e}")
        return text
    
    def _extract_from_txt(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading TXT: {e}")
            return ""
    
    def _extract_from_image(self, file_path: Path) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            print(f"Error performing OCR: {e}")
            return ""
    
    def create_notes(self, text: str, title: Optional[str] = None) -> Dict[str, any]:
        """
        Convert extracted text into structured notes.
        
        Args:
            text: Raw text content
            title: Optional title for the notes
            
        Returns:
            Dictionary containing structured notes
        """
        if not text or not text.strip():
            return {
                'title': title or "Empty Document",
                'sections': [],
                'summary': "No content to summarize."
            }
        
        # Clean and preprocess text
        cleaned_text = self._clean_text(text)
        
        # Extract key information
        sentences = sent_tokenize(cleaned_text)
        
        # Identify sections and structure
        sections = self._identify_sections(cleaned_text)
        
        # Generate summary
        summary = self._generate_summary(sentences)
        
        # Extract key concepts
        key_concepts = self._extract_key_concepts(cleaned_text)
        
        return {
            'title': title or self._generate_title(cleaned_text),
            'sections': sections,
            'key_concepts': key_concepts,
            'summary': summary
        }
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s.,!?;:()\-]', '', text)
        return text.strip()
    
    def _identify_sections(self, text: str) -> List[Dict[str, any]]:
        """Identify and structure sections in the text."""
        sections = []
        
        # Split by common section indicators
        lines = text.split('\n')
        current_section = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if line is a heading (short, capitalized, or numbered)
            if self._is_heading(line):
                # Save previous section
                if current_section:
                    sections.append({
                        'heading': current_section,
                        'content': self._format_content(current_content)
                    })
                current_section = line
                current_content = []
            else:
                current_content.append(line)
        
        # Add last section
        if current_section:
            sections.append({
                'heading': current_section,
                'content': self._format_content(current_content)
            })
        
        # If no sections found, create a single section
        if not sections:
            sections.append({
                'heading': 'Main Content',
                'content': self._format_content([text])
            })
        
        return sections
    
    def _is_heading(self, line: str) -> bool:
        """Determine if a line is likely a heading."""
        # Check various heading patterns
        if len(line) < 5 or len(line) > 100:
            return False
        
        # Numbered headings (1., 1.1, etc.)
        if re.match(r'^\d+\.(\d+\.)*\s+', line):
            return True
        
        # All caps or title case
        if line.isupper() or line.istitle():
            return True
        
        # Ends with colon
        if line.endswith(':'):
            return True
        
        return False
    
    def _format_content(self, content_lines: List[str]) -> List[str]:
        """Format content into bullet points."""
        formatted = []
        for line in content_lines:
            if line.strip():
                # Split long sentences
                sentences = sent_tokenize(line)
                for sentence in sentences:
                    if len(sentence.split()) > 5:  # Only add substantial sentences
                        formatted.append(sentence.strip())
        return formatted
    
    def _generate_summary(self, sentences: List[str], max_sentences: int = 3) -> str:
        """Generate a brief summary of the content."""
        if not sentences:
            return "No content to summarize."
        
        # Simple extractive summarization
        # Score sentences by word frequency
        word_freq = {}
        stop_words = set(stopwords.words('english'))
        
        for sentence in sentences:
            words = word_tokenize(sentence.lower())
            for word in words:
                if word.isalnum() and word not in stop_words:
                    word_freq[word] = word_freq.get(word, 0) + 1
        
        # Score sentences
        sentence_scores = {}
        for sentence in sentences:
            words = word_tokenize(sentence.lower())
            score = sum(word_freq.get(word, 0) for word in words if word.isalnum())
            sentence_scores[sentence] = score
        
        # Get top sentences
        top_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)
        summary_sentences = [s[0] for s in top_sentences[:max_sentences]]
        
        return ' '.join(summary_sentences)
    
    def _extract_key_concepts(self, text: str) -> List[str]:
        """Extract key concepts and terms from the text."""
        words = word_tokenize(text.lower())
        stop_words = set(stopwords.words('english'))
        
        # Filter and count words
        word_freq = {}
        for word in words:
            if word.isalnum() and len(word) > 3 and word not in stop_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top concepts
        top_concepts = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
        return [concept[0].capitalize() for concept in top_concepts]
    
    def _generate_title(self, text: str) -> str:
        """Generate a title from the content."""
        sentences = sent_tokenize(text)
        if sentences:
            # Use first sentence or extract key words
            first_sentence = sentences[0]
            words = first_sentence.split()[:6]
            return ' '.join(words) + ('...' if len(first_sentence.split()) > 6 else '')
        return "Lecture Notes"
    
    def export_to_txt(self, notes: Dict[str, any], output_path: str):
        """Export notes to TXT file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            # Title
            f.write(f"{'='*60}\n")
            f.write(f"{notes['title'].upper()}\n")
            f.write(f"{'='*60}\n\n")
            
            # Key Concepts
            if notes.get('key_concepts'):
                f.write("KEY CONCEPTS:\n")
                f.write(f"{'-'*60}\n")
                for concept in notes['key_concepts']:
                    f.write(f"• {concept}\n")
                f.write("\n")
            
            # Sections
            for section in notes['sections']:
                f.write(f"\n{section['heading'].upper()}\n")
                f.write(f"{'-'*60}\n")
                for point in section['content']:
                    f.write(f"• {point}\n")
                f.write("\n")
            
            # Summary
            if notes.get('summary'):
                f.write(f"\nSUMMARY\n")
                f.write(f"{'-'*60}\n")
                f.write(f"{notes['summary']}\n")
    
    def export_to_pdf(self, notes: Dict[str, any], output_path: str):
        """Export notes to PDF file."""
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Title
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, notes['title'], ln=True, align='C')
        pdf.ln(5)
        
        # Key Concepts
        if notes.get('key_concepts'):
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, 'KEY CONCEPTS:', ln=True)
            pdf.set_font('Arial', '', 10)
            for concept in notes['key_concepts']:
                pdf.multi_cell(0, 5, f'  • {concept}')
            pdf.ln(5)
        
        # Sections
        for section in notes['sections']:
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, section['heading'], ln=True)
            pdf.set_font('Arial', '', 10)
            for point in section['content']:
                # Handle long text
                pdf.multi_cell(0, 5, f'  • {point}')
            pdf.ln(3)
        
        # Summary
        if notes.get('summary'):
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, 'SUMMARY:', ln=True)
            pdf.set_font('Arial', '', 10)
            pdf.multi_cell(0, 5, notes['summary'])
        
        pdf.output(output_path)


def main():
    """Main function for command-line usage."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python bob_agent.py <input_file> [output_format]")
        print("Supported formats: pdf, ppt, pptx, txt, png, jpg, jpeg")
        print("Output formats: txt, pdf (default: both)")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_format = sys.argv[2] if len(sys.argv) > 2 else 'both'
    
    # Initialize Bob AI Agent
    bob = BobAIAgent()
    
    print(f"Processing file: {input_file}")
    
    # Process file
    extracted = bob.process_file(input_file)
    print(f"Extracted {len(extracted['text'])} characters")
    
    # Create notes
    notes = bob.create_notes(extracted['text'])
    print(f"Generated notes with {len(notes['sections'])} sections")
    
    # Export notes
    base_name = Path(input_file).stem
    
    if output_format in ['txt', 'both']:
        txt_output = f"{base_name}_notes.txt"
        bob.export_to_txt(notes, txt_output)
        print(f"Exported to: {txt_output}")
    
    if output_format in ['pdf', 'both']:
        pdf_output = f"{base_name}_notes.pdf"
        bob.export_to_pdf(notes, pdf_output)
        print(f"Exported to: {pdf_output}")
    
    print("\nNotes created successfully!")


if __name__ == "__main__":
    main()

# Made with Bob
